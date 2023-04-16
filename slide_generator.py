import io
import os
import random
import tempfile

import openai
import requests
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from utils import ask_chatgpt, get_dominant_colors, contrast_color, text_width, dim_image, adjust_contrast, apply_blur, \
    apply_overlay


def generate_title(prompt, max_tokens=40):
    return ask_chatgpt(prompt, max_tokens)


def generate_bullet_points(prompt, max_tokens=300):
    response = ask_chatgpt(prompt, max_tokens)

    bullet_points = response.strip().split("\n")
    return bullet_points


def add_picture_from_pil_image(slide, pil_image, left, top, width, height):
    # Save the PIL image to a temporary in-memory file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as image_file:
        pil_image.save(image_file, "PNG")
        image_file.seek(0)
        image_filename = image_file.name

    # Add the temporary file to the slide
    picture_shape = slide.Shapes.AddPicture(
        FileName=image_filename,
        LinkToFile=-1,  # Do not link to the file (embed the picture)
        SaveWithDocument=-1,  # Save the picture with the document
        Left=left,
        Top=top,
        Width=width,
        Height=height,
    )

    # Delete the temporary file from disk
    dominant_colors = get_dominant_colors(image_filename)
    text_color = contrast_color(dominant_colors[0])

    os.remove(image_filename)

    return text_color


def apply_background_manipulations(image_filename):
    # Load the image
    image = Image.open(image_filename)

    # Adjust brightness, contrast, blur, and overlay
    brightness_factor = 0.4
    contrast_factor = 3
    blur_radius = 2
    overlay_color = (0, 0, 0, 255)
    overlay_alpha = 0.3

    image = dim_image(image, brightness_factor)
    image = adjust_contrast(image, contrast_factor)
    image = apply_blur(image, blur_radius)
    image = apply_overlay(image, overlay_color, overlay_alpha)

    # Save the modified image
    image.save(image_filename)


def add_picture_from_pil_image_as_background(slide, presentation, pil_image):
    # Save the PIL image to a temporary in-memory file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as image_file:
        pil_image.save(image_file, "PNG")
        image_file.seek(0)
        image_filename = image_file.name

    apply_background_manipulations(image_filename)

    # Set the slide background image
    pic = slide.shapes.add_picture(image_filename, 0, 0, width=presentation.slide_width,
                                   height=presentation.slide_height)

    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    # Delete the temporary file from disk
    os.remove(image_filename)
    return


def add_related_picture(slide, presentation, prompt):
    # Generate the image using DALL-E
    response = openai.Image.create(prompt=prompt)

    # Get the image URL from the response
    image_url = response["data"][0]["url"]

    # Download the image from the URL
    image_data = requests.get(image_url).content

    # Load the image data into a PIL Image object
    pil_image = Image.open(io.BytesIO(image_data))

    add_picture_from_pil_image_as_background(slide, presentation, pil_image)
    return


def generate_slide(presentation, topic):
    prompt = f"Change the question: \"'{topic}'\" to a verbal noun title for an article. Start your title here:"
    title = generate_title(prompt).replace('"', '').replace('\n', '')
    print(f"Title: {title}")

    prompt = f"Please provide a summary and interesting information about the topic \"{title}\" using bullet points. Use the following format for your response:" \
             f"\n• Main Point 1\n--• Sub-point 1.1\n--• Sub-point 1.2\n• Main Point 2\n\nStart your response here:"
    bullet_points = generate_bullet_points(prompt)

    # Add a slide to the presentation
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    # Add background image that is related to the topic
    # add_related_picture(slide, presentation, title)

    # Set the slide title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Set the initial font size
    font_size = Pt(44)

    # Calculate the text width based on the initial font size
    text_width_pixels = text_width(title_shape.text, int(font_size))

    # Adjust the font size if the title exceeds the slide width
    slide_width_pixels = Inches(13)  # For 16:9

    while text_width_pixels > slide_width_pixels:
        font_size -= Pt(1)
        text_width_pixels = text_width(title_shape.text, int(font_size))

    # Apply the adjusted font size to the title
    title_shape.text_frame.paragraphs[0].runs[0].font.size = font_size
    title_shape.text_frame.paragraphs[0].runs[0].font.name = 'Goudy Old Style (Headings)'

    # Add bullet points to the slide
    left = Inches(1.3)
    top = Inches(1.2)
    width = Inches(random.randint(9, 10))  # Set the width to the random slide width
    height = Inches(6)

    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame

    # Adjust the text box properties
    tf.word_wrap = True  # Enable word wrapping
    tf.auto_size = 0  # Disable auto resizing
    tf.margin_left = Inches(0.1)  # Set left margin
    tf.margin_right = Inches(0.1)  # Set right margin
    tf.margin_top = Inches(0.1)  # Set top margin
    tf.margin_bottom = Inches(0.1)  # Set bottom margin

    # Add nested bullet points to the text of the presentation
    for point in bullet_points:
        if point.startswith("--"):
            p = tf.add_paragraph()
            p.text = point.strip("-")
            p.level = 1
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.name = "Avenir Next LT Pro (Body)"
            p.font.italic = True
            p.space_before = Pt(10)

        p.font.color.rgb = RGBColor(255, 255, 255)

    return
